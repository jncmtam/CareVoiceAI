#!/usr/bin/env python3
"""Generate CareVoice AI Vòng 2 MVP presentation for HackAIthon."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "/Users/jant/Desktop/CareVoiceAI_Vong2_MVP.pptx"

# Teal Trust + Sage Calm palette
TEAL = RGBColor(0x02, 0x80, 0x90)
SEAFOAM = RGBColor(0x00, 0xA8, 0x96)
MINT = RGBColor(0x02, 0xC3, 0x9A)
SAGE = RGBColor(0x84, 0xB5, 0x9F)
DARK = RGBColor(0x1E, 0x3A, 0x4A)
SLATE = RGBColor(0x50, 0x80, 0x8E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xFA, 0xF9)
LIGHT_TEAL = RGBColor(0xE0, 0xF4, 0xF1)
MUTED = RGBColor(0x5A, 0x6B, 0x73)
CORAL = RGBColor(0xE8, 0x5D, 0x4C)


def rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=13, color=DARK, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.font.bold = bold_first and i == 0
        p.space_after = Pt(6)
    return box


def slide_header(slide, title, subtitle=None):
    set_bg(slide, OFF_WHITE)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(5.625), TEAL)
    add_text_box(slide, Inches(0.55), Inches(0.35), Inches(8.8), Inches(0.7), title,
                 size=32, bold=True, color=DARK, font="Georgia")
    if subtitle:
        add_text_box(slide, Inches(0.55), Inches(0.95), Inches(8.5), Inches(0.45), subtitle,
                     size=14, color=MUTED)


def card(slide, x, y, w, h, title, body, accent=TEAL):
    add_rect(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, Inches(0.06), h, accent)
    add_text_box(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
                 title, size=15, bold=True, color=DARK)
    add_bullets(slide, x + Inches(0.18), y + Inches(0.48), w - Inches(0.3), h - Inches(0.55),
                body, size=11, color=MUTED)


def stat_callout(slide, x, y, number, label, accent=MINT):
    add_rect(slide, x, y, Inches(2.05), Inches(1.35), WHITE)
    add_rect(slide, x, y, Inches(2.05), Inches(0.08), accent)
    add_text_box(slide, x + Inches(0.12), y + Inches(0.2), Inches(1.8), Inches(0.65),
                 number, size=34, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.85), Inches(1.85), Inches(0.4),
                 label, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ── 1. Cover ──────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    add_rect(s, Inches(0), Inches(0), Inches(10), Inches(0.18), MINT)
    add_rect(s, Inches(7.2), Inches(1.2), Inches(2.5), Inches(2.5), TEAL)
    add_rect(s, Inches(7.5), Inches(1.5), Inches(1.9), Inches(1.9), SEAFOAM)
    add_text_box(s, Inches(0.7), Inches(1.4), Inches(6.2), Inches(1.1),
                 "CareVoice AI", size=48, bold=True, color=WHITE, font="Georgia")
    add_text_box(s, Inches(0.7), Inches(2.45), Inches(6.2), Inches(0.55),
                 "Điều dưỡng ảo 24/7 — MVP Vòng 2", size=22, color=MINT)
    add_text_box(s, Inches(0.7), Inches(3.2), Inches(6.5), Inches(0.9),
                 "HackAIthon VNPT · Đề tài 4 · Bảng B Challenger\n"
                 "Đội TMT · Đại học Bách Khoa TP.HCM",
                 size=14, color=RGBColor(0xB8, 0xD4, 0xD8))
    add_text_box(s, Inches(0.7), Inches(4.55), Inches(6), Inches(0.4),
                 "Voice-first · Explainable AI · Connected Care", size=12, color=SAGE)

    # ── 2. Vòng 1 → Vòng 2 ───────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Từ đề xuất Vòng 1 đến MVP Vòng 2", "Phát triển liên tục, không pivot")
    card(s, Inches(0.55), Inches(1.45), Inches(4.2), Inches(3.55),
         "Vòng 1 — Đề xuất",
         ["Voice-first outbound: AI chủ động hỏi thăm mỗi sáng",
          "Tích hợp VNPT: SmartReader, SmartVoice, SmartBot",
          "Mô hình B2B SaaS — phòng khám nội tiết & tim mạch",
          "Pilot TP.HCM & Hà Nội · TAM/SAM/SOM đã xác định"],
         accent=SAGE)
    card(s, Inches(5.0), Inches(1.45), Inches(4.45), Inches(3.55),
         "Vòng 2 — MVP đã code",
         ["iOS native SwiftUI (79 file) + FastAPI backend",
          "Dashboard điều dưỡng ưu tiên + gọi 1 chạm",
          "Buổi sáng 2 bước: thuốc → lời khuyên",
          "20 pytest · Docker · offline queue · job runner"],
         accent=MINT)

    # ── 3. Vấn đề & Giải pháp ────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Vấn đề & Giải pháp", "Nhân văn làm lõi, công nghệ làm cầu nối")
    add_rect(s, Inches(0.55), Inches(1.5), Inches(4.15), Inches(3.4), RGBColor(0xFD, 0xED, 0xEC))
    add_rect(s, Inches(0.55), Inches(1.5), Inches(4.15), Inches(0.5), CORAL)
    add_text_box(s, Inches(0.75), Inches(1.58), Inches(3.8), Inches(0.4),
                 "Pain points", size=16, bold=True, color=WHITE)
    add_bullets(s, Inches(0.75), Inches(2.15), Inches(3.75), Inches(2.6),
                ["BN cao tuổi: chữ nhỏ, không rành smartphone",
                 "Điều dưỡng: 8–10 phút/cuộc gọi check-in, thiếu nhân lực",
                 "Người nhà: biết muộn khi có dấu hiệu bất thường",
                 "Triệu chứng nhẹ dễ lọt · nặng đôi khi đến muộn"],
                size=12, color=DARK)
    add_rect(s, Inches(5.0), Inches(1.5), Inches(4.45), Inches(3.4), LIGHT_TEAL)
    add_rect(s, Inches(5.0), Inches(1.5), Inches(4.45), Inches(0.5), TEAL)
    add_text_box(s, Inches(5.2), Inches(1.58), Inches(4), Inches(0.4),
                 "CareVoice AI", size=16, bold=True, color=WHITE)
    add_bullets(s, Inches(5.2), Inches(2.15), Inches(4.1), Inches(2.6),
                ["Voice-first: TTS đọc câu hỏi, STT nhận giọng nói",
                 "Explainable triage: lý do phân loại minh bạch",
                 "SMS người nhà tự động khi cần chú ý/can thiệp",
                 "Dashboard ưu tiên + gọi 1 chạm cho điều dưỡng"],
                size=12, color=DARK)

    # ── 4. Link sản phẩm / MVP ───────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Link sản phẩm & MVP", "Có thể thao tác và kiểm thử ngay")
    card(s, Inches(0.55), Inches(1.45), Inches(2.85), Inches(2.2),
         "iOS App (SwiftUI)",
         ["Repo: HackAIthon/CareVoiceAI",
          "Demo mode: không cần backend",
          "Xcode 15+ · iOS 17+"],
         accent=TEAL)
    card(s, Inches(3.55), Inches(1.45), Inches(2.85), Inches(2.2),
         "Backend API",
         ["FastAPI · OpenAPI docs",
          "http://127.0.0.1:8000/api/v1",
          "Docker Compose sẵn sàng"],
         accent=SEAFOAM)
    card(s, Inches(6.55), Inches(1.45), Inches(2.85), Inches(2.2),
         "Tài khoản demo",
         ["Staff: nurse01@hospital.vn",
          "BN: BN-2026-0001 / 4567",
          "OTP demo: 123456"],
         accent=MINT)
    add_rect(s, Inches(0.55), Inches(3.85), Inches(8.85), Inches(1.35), WHITE)
    add_text_box(s, Inches(0.75), Inches(3.95), Inches(8.5), Inches(1.1),
                 "MVP trực quan: ứng dụng iOS native + REST API + dashboard điều dưỡng\n"
                 "Tài liệu: APP_FEATURES_AND_FLOWS.md · API_CONTRACT.md · backend/README.md",
                 size=12, color=MUTED)

    # ── 5. Kiến trúc ───────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Kiến trúc hệ thống", "App không gọi trực tiếp VNPT — backend làm gateway")
    layers = [
        ("iOS App", "SwiftUI · Demo/Backend mode\nOffline queue · Local notifications", TEAL),
        ("CareVoice API", "FastAPI · JWT · RBAC · Job runner\nPostgreSQL/SQLite · Redis", SEAFOAM),
        ("VNPT Gateway", "SmartReader · SmartVoice STT/TTS\nSmartBot mock/live", MINT),
    ]
    for i, (title, body, accent) in enumerate(layers):
        y = Inches(1.45 + i * 1.15)
        add_rect(s, Inches(0.55), y, Inches(8.85), Inches(1.0), WHITE)
        add_rect(s, Inches(0.55), y, Inches(1.6), Inches(1.0), accent)
        add_text_box(s, Inches(0.65), y + Inches(0.28), Inches(1.4), Inches(0.45),
                     title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(2.35), y + Inches(0.18), Inches(6.8), Inches(0.7),
                     body, size=12, color=MUTED)
    add_text_box(s, Inches(0.55), Inches(4.85), Inches(8.5), Inches(0.4),
                 "Luồng: BN gửi check-in → STT + phân loại → SMS người nhà + cảnh báo dashboard",
                 size=11, color=SLATE, align=PP_ALIGN.CENTER)

    # ── 6. Tính năng BN ────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Tính năng đã triển khai — Bệnh nhân", "5 tab · Voice-first · Buổi sáng 2 bước")
    features = [
        ("Check-in hàng ngày", "TTS câu hỏi · nút Có/Không/Bình thường\nGhi âm STT · polling kết quả · offline queue"),
        ("Thuốc & tuân thủ", "Nhắc giọng theo khung giờ · xác nhận buổi\nPOST adherence · điều dưỡng thấy liều bỏ lỡ"),
        ("Hotline AI 24/7", "Hỏi bằng chữ/giọng · guardrail thuốc\nneeds_staff_review khi cần người thật"),
        ("Lời khuyên hôm nay", "Mẹo sức khỏe theo chẩn đoán\nTick bước 2 buổi sáng"),
    ]
    for i, (t, b) in enumerate(features):
        col, row = i % 2, i // 2
        x = Inches(0.55 + col * 4.65)
        y = Inches(1.45 + row * 1.85)
        card(s, x, y, Inches(4.35), Inches(1.65), t, [b], accent=TEAL if i % 2 == 0 else SAGE)

    # ── 7. Tính năng Staff ─────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Tính năng đã triển khai — Điều dưỡng", "Ưu tiên đúng người · ít gõ phím")
    card(s, Inches(0.55), Inches(1.45), Inches(4.2), Inches(1.75),
         "Dashboard ưu tiên",
         ["Sắp xếp theo mức nguy cơ: normal / attention / intervention",
          "KPI: phút tiết kiệm, check-in hoàn thành, danh sách ưu tiên",
          "Critical alert banner + haptics khi có ca mới"],
         accent=CORAL)
    card(s, Inches(5.0), Inches(1.45), Inches(4.45), Inches(1.75),
         "Hồ sơ & timeline",
         ["Timeline minh bạch: lý do phân loại AI",
          "Gọi 1 chạm bệnh nhân & người nhà (tel:)",
          "Cập nhật trạng thái xử lý handling"],
         accent=TEAL)
    card(s, Inches(0.55), Inches(3.45), Inches(4.2), Inches(1.75),
         "OCR đơn thuốc",
         ["Upload PDF/ảnh → job polling",
          "SmartReader scan → needs_review → confirm",
          "Tự động tạo lịch nhắc thuốc"],
         accent=SEAFOAM)
    card(s, Inches(5.0), Inches(3.45), Inches(4.45), Inches(1.75),
         "Tạo BN mới",
         ["Validation mã BN + SĐT chuẩn hóa",
          "Nhập thuốc, lịch tái khám",
          "RBAC: nurse/doctor/admin"],
         accent=MINT)

    # ── 8. API VNPT ────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "API VNPT đã tích hợp", "Backend gateway — iOS không lộ credential")
    rows = [
        ["VNPT API", "Use case CareVoice", "Endpoint CareVoice"],
        ["SmartReader OCR", "Scan đơn thuốc, giấy ra viện", "POST /patients/{id}/documents"],
        ["SmartVoice STT", "Nhận giọng check-in & hotline", "POST /checkins/{id}/responses"],
        ["SmartVoice TTS", "Đọc câu hỏi check-in", "GET /checkins/{id}/audio"],
        ["SmartBot / Tóm tắt", "Phân loại nguy cơ + hotline", "GET /checkin_jobs/{job_id}"],
    ]
    tbl = s.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.4), Inches(8.85), Inches(3.5)).table
    col_w = [Inches(2.2), Inches(3.3), Inches(3.35)]
    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if ri > 0 else 12)
                p.font.name = "Calibri"
                p.font.bold = ri == 0
                p.font.color.rgb = WHITE if ri == 0 else DARK
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TEAL
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_TEAL
    add_text_box(s, Inches(0.55), Inches(5.05), Inches(8.5), Inches(0.35),
                 "Mock adapter + live gateway tại app/integrations/vnpt/ · 7 test STT parser",
                 size=10, color=MUTED)

    # ── 9. Demo flow ───────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Flow demo 90 giây", "Buổi sáng của bác Trần Văn Bình")
    steps = [
        ("1", "Đăng nhập BN", "BN-2026-0001 · nghe lời chào buổi sáng"),
        ("2", "Check-in", "Chọn Có → chóng mặt → AI phân loại + SMS người nhà"),
        ("3", "Dashboard", "Điều dưỡng rung · banner đỏ · gọi 1 chạm"),
        ("4", "Thuốc", "Xác nhận buổi sáng · tick 3/3"),
        ("5", "Hotline", "Hỏi quên liều · AI thận trọng + báo staff"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.4 + i * 0.78)
        add_rect(s, Inches(0.55), y, Inches(0.55), Inches(0.55), TEAL if i % 2 == 0 else SEAFOAM)
        add_text_box(s, Inches(0.55), y + Inches(0.08), Inches(0.55), Inches(0.4),
                     num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.3), y + Inches(0.05), Inches(2.5), Inches(0.35),
                     title, size=14, bold=True, color=DARK)
        add_text_box(s, Inches(3.9), y + Inches(0.05), Inches(5.4), Inches(0.45),
                     desc, size=12, color=MUTED)
    add_rect(s, Inches(0.55), Inches(5.15), Inches(8.85), Inches(0.35), MINT)
    add_text_box(s, Inches(0.7), Inches(5.18), Inches(8.5), Inches(0.3),
                 "Script thuyết trình Nam+Nữ ~4:30 trong APP_FEATURES_AND_FLOWS.md §13",
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 10. Video demo ─────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Video demo", "Quay trực tiếp test sản phẩm trên thiết bị thật")
    add_rect(s, Inches(0.55), Inches(1.45), Inches(5.5), Inches(3.5), DARK)
    add_rect(s, Inches(0.85), Inches(1.75), Inches(4.9), Inches(2.9), SLATE)
    add_text_box(s, Inches(1.2), Inches(2.6), Inches(4.2), Inches(1.2),
                 "▶  VIDEO DEMO\n\n"
                 "[Điền link YouTube/Drive\n khi upload video]",
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(6.3), Inches(1.55), Inches(3.2), Inches(3.3),
                ["Cảnh 1: BN check-in bằng giọng (30s)",
                 "Cảnh 2: Banner SMS người nhà (15s)",
                 "Cảnh 3: Dashboard điều dưỡng + gọi 1 chạm (25s)",
                 "Cảnh 4: Nhắc thuốc + hotline (20s)",
                 "Cảnh 5: OCR đơn thuốc (20s)"],
                size=12, color=DARK)
    add_text_box(s, Inches(0.55), Inches(5.05), Inches(8.5), Inches(0.35),
                 "Gợi ý: quay 2 iPhone — Lan (BN) + Minh (staff) theo script §13",
                 size=10, color=MUTED)

    # ── 11. Cài đặt iOS ────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Hướng dẫn cài đặt — iOS", "Demo mode hoạt động ngay không cần server")
    add_rect(s, Inches(0.55), Inches(1.45), Inches(8.85), Inches(3.55), RGBColor(0x2D, 0x2D, 0x2D))
    code = (
        "# Clone repo\n"
        "git clone <repo-url> && cd HackAIthon\n\n"
        "# Mở Xcode\n"
        "open CareVoiceAI.xcodeproj\n\n"
        "# Chọn simulator hoặc thiết bị thật → Run (⌘R)\n"
        "# Demo mode ON (mặc định) — đăng nhập ngay\n\n"
        "# Kết nối backend thật:\n"
        "# Settings → tắt Demo mode → Base URL http://127.0.0.1:8000/api/v1"
    )
    add_text_box(s, Inches(0.75), Inches(1.6), Inches(8.5), Inches(3.2),
                 code, size=11, color=RGBColor(0xA8, 0xE6, 0xCF), font="Consolas")

    # ── 12. Cài đặt Backend ────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Hướng dẫn cài đặt — Backend", "Local Python hoặc Docker Compose")
    card(s, Inches(0.55), Inches(1.45), Inches(4.2), Inches(3.55),
         "Chạy local",
         ["cd backend",
          "python3.12 -m venv .venv",
          ".venv/bin/pip install -e \".[dev]\"",
          "uvicorn app.main:app --reload --port 8000",
          "OpenAPI: /api/v1/docs"],
         accent=TEAL)
    card(s, Inches(5.0), Inches(1.45), Inches(4.45), Inches(3.55),
         "Docker Compose",
         ["cd backend && docker compose up --build",
          "API + PostgreSQL + Redis",
          "Seed demo accounts tự động",
          "VENDOR_MOCK_MODE=true (mặc định demo)"],
         accent=SEAFOAM)

    # ── 13. Script test tự động ────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Script test tự động", "20 pytest — contract, jobs, storage, STT")
    stat_callout(s, Inches(0.55), Inches(1.5), "20", "tests passed", MINT)
    stat_callout(s, Inches(2.75), Inches(1.5), "5", "test modules", TEAL)
    stat_callout(s, Inches(4.95), Inches(1.5), "<1s", "collect time", SEAFOAM)
    add_rect(s, Inches(0.55), Inches(3.1), Inches(8.85), Inches(2.0), RGBColor(0x2D, 0x2D, 0x2D))
    test_code = (
        "cd backend\n"
        ".venv/bin/python -m pytest -v\n\n"
        "# Modules:\n"
        "# test_api_contract.py — auth, dashboard, check-in, OCR\n"
        "# test_background_jobs.py — job_runner async\n"
        "# test_storage_limits.py — 413/415 upload limits\n"
        "# test_vnpt_stt.py — STT parser & endpoints\n"
        "# test_patient_validation.py — BN code normalization"
    )
    add_text_box(s, Inches(0.75), Inches(3.25), Inches(8.5), Inches(1.7),
                 test_code, size=11, color=RGBColor(0xA8, 0xE6, 0xCF), font="Consolas")

    # ── 14. Phân khúc & doanh thu ───────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Chiến lược triển khai — Thị trường", "B2B SaaS · Phòng khám chuyên khoa mạn tính")
    stat_callout(s, Inches(0.55), Inches(1.5), "~8M", "TAM BN mạn tính VN", TEAL)
    stat_callout(s, Inches(2.75), Inches(1.5), "~400K", "SAM quản lý tại PK", SEAFOAM)
    stat_callout(s, Inches(4.95), Inches(1.5), "2–5K", "SOM pilot 12 tháng", MINT)
    card(s, Inches(0.55), Inches(3.05), Inches(4.2), Inches(2.1),
         "Phân khúc khách hàng",
         ["Phòng khám nội tiết & tim mạch (TP.HCM, Hà Nội)",
          "BN đái tháo đường, THA, suy tim — sống tại nhà",
          "Điều dưỡng chăm sóc tại nhà & telecare"],
         accent=SAGE)
    card(s, Inches(5.0), Inches(3.05), Inches(4.45), Inches(2.1),
         "Mô hình doanh thu",
         ["SaaS theo số BN active/tháng",
          "Gói Professional: dashboard + OCR + voice",
          "Add-on: SMS gateway, tích hợp HIS"],
         accent=TEAL)

    # ── 15. Pricing & unit economics ───────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Định giá & Unit Economics", "Chi phí biến đổi theo request VNPT")
    rows = [
        ["Hạng mục", "Đơn giá ước tính", "Ghi chú"],
        ["SaaS / BN / tháng", "45.000 – 80.000 VNĐ", "Theo gói & volume"],
        ["STT request", "~800 VNĐ / phút audio", "SmartVoice async"],
        ["TTS request", "~500 VNĐ / 1K ký tự", "Câu hỏi check-in cache"],
        ["OCR page", "~1.200 VNĐ / trang", "SmartReader scan"],
        ["Gross margin mục tiêu", "55–65%", "Sau pilot 500 BN"],
    ]
    tbl = s.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.4), Inches(8.85), Inches(2.8)).table
    for ci, w in enumerate([Inches(2.8), Inches(2.8), Inches(3.25)]):
        tbl.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Calibri"
                p.font.bold = ri == 0
                p.font.color.rgb = WHITE if ri == 0 else DARK
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SEAFOAM
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_TEAL
    add_bullets(s, Inches(0.55), Inches(4.4), Inches(8.85), Inches(1.0),
                ["Pilot 3 phòng khám × 150 BN = 450 BN → MRR ~27–36 triệu VNĐ",
                 "CAC qua đối tác phòng khám · LTV/CAC mục tiêu > 4× trong 18 tháng"],
                size=12, color=DARK)

    # ── 16. GTM & đối tác ──────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "GTM · Kênh phân phối · Đối tác", "Land-and-expand từ pilot")
    channels = [
        ("Pilot PK", "3 PK nội tiết/tim mạch\nTP.HCM & Hà Nội · 90 ngày", TEAL),
        ("Đối tác VNPT", "SmartVoice/SmartReader\nbundle API + hỗ trợ kỹ thuật", SEAFOAM),
        ("Bệnh viện công", "Telecare post-discharge\nliên kết khoa nội trú", MINT),
        ("Truyền thông", "Case study người nhà\nDemo video · hội thảo y tế", SAGE),
    ]
    for i, (t, b, accent) in enumerate(channels):
        col, row = i % 2, i // 2
        card(s, Inches(0.55 + col * 4.65), Inches(1.45 + row * 1.95),
             Inches(4.35), Inches(1.75), t, [b], accent=accent)

    # ── 17. Roadmap 12 tháng ───────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "Lộ trình mở rộng 12 tháng", "Từ MVP hackathon → production pilot")
    roadmap = [
        ("Q3/26", "MVP Vòng 2", "iOS + API + VNPT mock/live\n20 test · Docker"),
        ("Q4/26", "Pilot 1", "1 PK · 150 BN · SMS thật\nMetrics dashboard"),
        ("Q1/27", "Scale pilot", "3 PK · HIS integration POC\nMetrics & reporting"),
        ("Q2/27", "Commercial", "Pricing launch · 500 BN\nAPNs push · Redis queue"),
    ]
    add_rect(s, Inches(0.55), Inches(2.0), Inches(8.85), Inches(0.08), MINT)
    for i, (q, title, body) in enumerate(roadmap):
        x = Inches(0.55 + i * 2.25)
        add_rect(s, x + Inches(0.75), Inches(1.85), Inches(0.2), Inches(0.2), TEAL)
        add_text_box(s, x, Inches(1.35), Inches(2.1), Inches(0.4), q, size=11, bold=True, color=TEAL)
        add_text_box(s, x, Inches(2.2), Inches(2.1), Inches(0.4), title, size=13, bold=True, color=DARK)
        add_text_box(s, x, Inches(2.65), Inches(2.1), Inches(1.5), body, size=10, color=MUTED)
    add_rect(s, Inches(0.55), Inches(4.35), Inches(8.85), Inches(1.0), LIGHT_TEAL)
    add_text_box(s, Inches(0.75), Inches(4.5), Inches(8.5), Inches(0.7),
                 "Scale kỹ thuật: VNPT live gateway · Alembic migration · S3 storage · "
                 "Celery/RQ jobs · Prometheus metrics · retention PHI/PII",
                 size=11, color=DARK)

    # ── 18. KPI & Impact ───────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(s, "KPI pilot & Tác động", "Đo bằng phút tiết kiệm — và sự yên tâm của người nhà")
    stat_callout(s, Inches(0.55), Inches(1.5), "70%", "check-in tự phục vụ", MINT)
    stat_callout(s, Inches(2.75), Inches(1.5), "8→2", "phút/cuộc gọi", TEAL)
    stat_callout(s, Inches(4.95), Inches(1.5), "<5'", "phản hồi ca attention", SEAFOAM)
    stat_callout(s, Inches(7.15), Inches(1.5), "85%", "tuân thủ thuốc", SAGE)
    card(s, Inches(0.55), Inches(3.05), Inches(4.2), Inches(2.1),
         "Tác động xã hội",
         ["Giảm tải điều dưỡng — tập trung ca cần can thiệp",
          "Người nhà được báo sớm, không phụ thuộc con cái đi làm",
          "BN cao tuổi: 30 giây giọng nói thay vì app phức tạp"],
         accent=TEAL)
    card(s, Inches(5.0), Inches(3.05), Inches(4.45), Inches(2.1),
         "Khác biệt cốt lõi",
         ["Voice-first — không phải chatbot đa năng",
          "Explainable AI — điều dưỡng tin vì thấy lý do",
          "Offline-first — người già không chịu lỗi Wi-Fi"],
         accent=MINT)

    # ── 19. Q&A / Thank you ─────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    add_rect(s, Inches(0), Inches(5.445), Inches(10), Inches(0.18), MINT)
    add_text_box(s, Inches(0.7), Inches(1.5), Inches(8.5), Inches(1.0),
                 "Cảm ơn & Q&A", size=44, bold=True, color=WHITE, font="Georgia",
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(2.7), Inches(8.5), Inches(1.2),
                 "CareVoice AI — Mỗi buổi sáng, một câu hỏi.\n"
                 "Mỗi câu trả lời, đúng người được gọi.",
                 size=20, color=MINT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(4.0), Inches(8.5), Inches(1.0),
                 "Đội TMT · ĐH Bách Khoa TP.HCM\n"
                 "Repo: HackAIthon · Demo sẵn sàng trên iOS & API",
                 size=14, color=RGBColor(0xB8, 0xD4, 0xD8), align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()